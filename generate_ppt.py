import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx is not installed. Installing it now...")
    import subprocess
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except Exception as e:
        print(f"Error installing python-pptx: {e}")
        print("Please run: pip install python-pptx manually in your terminal.")
        sys.exit(1)

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define colors matching Groundwork theme
    c_dark = RGBColor(8, 11, 17)       # Deep Space Navy
    c_white = RGBColor(243, 244, 246)  # Main Text
    c_muted = RGBColor(156, 163, 175)  # Muted Text
    c_violet = RGBColor(139, 92, 246)  # Primary Purple
    c_cyan = RGBColor(6, 182, 212)     # Accent Cyan
    
    # ----------------------------------------------------
    # Slide 1: Title Slide
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c_dark
    
    # Title Text Box
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "GROUNDWORK"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = c_violet
    p.font.name = "Space Grotesk"
    p.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "ML-Powered Startup Growth & Investment Predictor"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = c_cyan
    p2.font.name = "Plus Jakarta Sans"
    
    # Subtitle Text Box
    txBox_sub = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.0))
    tf_sub = txBox_sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "A Data-Driven Learning & Vetting Platform for Angel Investors"
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = c_muted
    p_sub.font.name = "Plus Jakarta Sans"
    
    # ----------------------------------------------------
    # Slide Helper for Content Slides
    # ----------------------------------------------------
    def add_content_slide(slide_num, title_text, points):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = c_dark
        
        # Slide Number & Title
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.3), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p_num = tf.paragraphs[0]
        p_num.text = f"0{slide_num} // GROUNDWORK PROJECT"
        p_num.font.size = Pt(11)
        p_num.font.bold = True
        p_num.font.color.rgb = c_muted
        p_num.font.name = "Space Grotesk"
        p_num.space_after = Pt(4)
        
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = c_violet
        p_title.font.name = "Space Grotesk"
        
        # Content Box (Two-column layout if points length > 2)
        txBox_content = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(4.5))
        tf_content = txBox_content.text_frame
        tf_content.word_wrap = True
        
        for idx, (headline, body) in enumerate(points):
            p_pt = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
            p_pt.space_after = Pt(16)
            
            run_head = p_pt.add_run()
            run_head.text = f"•  {headline}  |  "
            run_head.font.bold = True
            run_head.font.size = Pt(16)
            run_head.font.color.rgb = c_cyan
            run_head.font.name = "Plus Jakarta Sans"
            
            run_body = p_pt.add_run()
            run_body.text = body
            run_body.font.size = Pt(15)
            run_body.font.color.rgb = c_white
            run_body.font.name = "Plus Jakarta Sans"
            
        return slide

    # ----------------------------------------------------
    # Slide 2: The Problem
    # ----------------------------------------------------
    add_content_slide(2, "The Problem Statement", [
        ("High Early-Stage Risk", "Up to 90% of early-stage startups fail. New angel investors enter the ecosystem with high capital risk and minimal data support."),
        ("Information Asymmetry", "Novice investors struggle to audit operational health, verify customer unit economics, and distinguish solid metrics from marketing hype."),
        ("Jargon Barrier", "Venture capital concepts like TAM, LTV, Churn, and Runway create an entry barrier for beginners trying to analyze financial metrics."),
        ("No Sandboxes", "No secure educational environments exist where aspiring investors can practice evaluating real-world business outcomes without risking capital.")
    ])

    # ----------------------------------------------------
    # Slide 3: The Groundwork Solution
    # ----------------------------------------------------
    add_content_slide(3, "The Groundwork Solution", [
        ("ML-Driven Screening", "Replaces gut-feel investing with empirical modeling, evaluating startup profiles against 1,000 synthesized business records."),
        ("Beginner-First Education", "Provides on-the-fly hover tooltips for all prediction forms, translating professional finance terms into accessible analogies."),
        ("Interactive Funding Tutorials", "Provides structured explanations on the progression and relationship between Pre-seed, Seed, Series A, and Series B stages."),
        ("Practice Sandbox Arena", "Includes an interactive venture game with 6 business scenarios to build core metrics-appraisal instincts risk-free.")
    ])

    # ----------------------------------------------------
    # Slide 4: Core Dashboard Features
    # ----------------------------------------------------
    add_content_slide(4, "Core Dashboard Modules", [
        ("Overview Console", "A central landing showing screen sizes, top investment picks, and 36-month line trajectories comparing success vs. failure index benchmarks."),
        ("Startup Explorer", "A fully filterable, searchable directory grid. Users can click 'Analyze' on any record to view normalized metrics compared to cohort averages."),
        ("Prediction Engine", "An interactive calculator allowing investors to input custom startup statistics to retrieve dynamic growth classes and ROI estimations."),
        ("Practice Simulator", "A 6-level interactive quiz displaying expert rationale and color-coded scorecards showing 'CORRECT' or 'WRONG' answers.")
    ])

    # ----------------------------------------------------
    # Slide 5: Data Synthesis & Metrics
    # ----------------------------------------------------
    add_content_slide(5, "Cohort Features & Dataset", [
        ("Synthesized Cohorts", "Simulates 1,000 detailed startup profiles across 7 industry sectors (AI, SaaS, FinTech, HealthTech, CleanTech, E-Commerce, EdTech)."),
        ("Unit Economics Focus", "Records customer metrics like Customer Lifetime Value (LTV), Customer Acquisition Cost (CAC), and LTV/CAC ratios."),
        ("Operational Parameters", "Collects Customer Retention/Churn, Net Promoter Score (NPS), and Website Traffic Growth."),
        ("Capital Runway Indicators", "Calculates capital runway based on cash reserves vs. monthly cash expenditure (Burn Rate).")
    ])

    # ----------------------------------------------------
    # Slide 6: Machine Learning Architecture
    # ----------------------------------------------------
    add_content_slide(6, "Predictive Modeling Core", [
        ("Dual-Model Pipeline", "Trains two separate Random Forest models built on standardized numerical scaling and one-hot categorical encoding."),
        ("Classification Engine", "Categorizes startups into Low, Medium, and High Growth Potential tiers using Random Forest Classifier probabilities."),
        ("Regression Engine", "Estimates the expected 3-year valuation multiple (ROI) using a skewed power-law distribution typical of venture capital returns."),
        ("Evaluation Results", "Groundwork achieves ~84% accuracy in classification, and captures over 76% of ROI variances (R² score) in the cohort.")
    ])

    # ----------------------------------------------------
    # Slide 7: Interactive Practice Arena
    # ----------------------------------------------------
    add_content_slide(7, "The Venture Practice Arena", [
        ("Gamified Evaluation", "Offers 6 distinct operational case studies simulating common startup profiles (FinTech, EdTech, CleanTech, SaaS)."),
        ("Asymmetric Metrics", "Presents complex cases (e.g., hyper-growth SaaS with low customer efficiency, or R&D deeptech with high burn)."),
        ("Instant Decision Assessment", "Tells the user whether their guess (Invest/Decline) is correct or incorrect relative to model logic."),
        ("Expert Rationale", "Shows a bold feedback banner with the correct option and a full unit-economic analysis explaining the outcome.")
    ])

    # ----------------------------------------------------
    # Slide 8: Technical Jargon & Glossary
    # ----------------------------------------------------
    add_content_slide(8, "Educational & Accessibility Focus", [
        ("Hover Glossary Tooltips", "Interactive (i) icons explaining every form metric (NPS, Churn, TAM, CAC) upon cursor hover."),
        ("Real-World Analogies", "Simplifies complex terms (e.g., Churn: 'water leaking out of the bucket'; TAM: 'total size of the market cake')."),
        ("Moving Ticker Bar", "A rolling news ticker at the top displays index trends and sector multiples (e.g., SaaS average multiples)."),
        ("Technical Glossary", "A centralized, toggleable directory detailing VC terminology, formal definitions, and formulas.")
    ])

    # ----------------------------------------------------
    # Slide 9: Tech Stack & System Design
    # ----------------------------------------------------
    add_content_slide(9, "Tech Stack & System Design", [
        ("FastAPI Backend", "Exposes API endpoints (`/api/startups`, `/api/metrics`, `/api/predict`) and serves static resources using Python."),
        ("Vanilla Frontend", "Sleek single-page interface using HTML5, Vanilla CSS3, and JavaScript, optimized for fast loading and responsive layout."),
        ("Rich Glassmorphism UI", "Employs high-end dark backgrounds, border glows, sliding transitions, and backdrop blur filters."),
        ("Chart.js Visualization", "Renders dynamic sector doughnuts, model diagnostics (feature importances), and benchmark comparative charts.")
    ])

    # ----------------------------------------------------
    # Slide 10: Value & Next Milestones
    # ----------------------------------------------------
    add_content_slide(10, "Impact & Future Milestones", [
        ("Empowering Angel Investors", "Provides new angels with objective tools to screen companies, avoiding early-stage cash-burn traps."),
        ("Simplified Deployment", "Pre-configured requirements ready for standard container hosting on cloud platforms like Render or Railway."),
        ("Equity API Connection", "Future integrations with Crunchbase, AngelList, or SEMRush APIs to pull live startup indexes."),
        ("Community Leaderboard", "Planned features for group syndicates, collaborative vetting practice, and cohort scoring tables.")
    ])

    # Save presentation
    output_dir = "C:/Users/Lenovo/.gemini/antigravity/scratch/Project"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "Groundwork_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully at {output_path}")

if __name__ == "__main__":
    create_presentation()
